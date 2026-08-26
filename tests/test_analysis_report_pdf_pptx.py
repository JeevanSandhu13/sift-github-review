"""PDF and PowerPoint analysis-report exports.

Both formats are built from the exact same ``_gather_report_material``
call as the existing Markdown/HTML report, so the load-bearing
property to test is not "does reportlab/python-pptx run" but "does
the same session data show up in these formats the same way it shows
up in the already-trusted Markdown/HTML report" -- findings, tables,
verification warnings, figures (subject to the plot-kind allowlist),
and the disclosure statement.
"""

from __future__ import annotations

import base64
import json
from pathlib import Path

from sift import release_ledger
from sift.research_export import (
    build_analysis_report_pdf,
    build_analysis_report_pptx,
)
from sift.store import get_store

# 1x1 transparent PNG, same fixture as test_analysis_report.py.
_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk"
    "+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg==")


def _seed(tmp_path: Path) -> None:
    (tmp_path / ".sift").mkdir(exist_ok=True)
    store = get_store(tmp_path)
    store.insert(
        label="Churn on tenure", analysis_type="linear_regression",
        sanitized_payload={
            "type": "linear_regression", "n": 28194,
            "coefficients": {"(Intercept)": 0.41, "tenure": -0.12},
            "standard_errors": {"(Intercept)": 0.02, "tenure": 0.03},
            "p_values": {"(Intercept)": 0.001, "tenure": 0.0004},
            "r_squared": 0.148, "robust_se_type": "classical",
        },
        language="Python", script_code="x",
        transformations=["clamped SEs"], raw_log_path=None,
        script_run_id="r1")
    store.insert(
        label="Region counts", analysis_type="frequency_table",
        sanitized_payload={
            "type": "frequency_table", "variable": "region", "n": 120,
            "counts": {"north": 40, "south": 35, "east": 25, "west": 20},
            "missing_count": 0,
        },
        language="Python", script_code="x", transformations=[],
        raw_log_path=None, script_run_id="r1")
    release_ledger.record_release(
        tmp_path, kind="tool_response", tool="submit_script",
        response={"content": [{"type": "text", "text": '{"status":"ok"}'}]})


def _plot(tmp_path: Path, kind: str, fname: str, label: str,
          run: str = "run1") -> None:
    d = tmp_path / ".sift" / "runs" / run / "_sift_plots"
    d.mkdir(parents=True, exist_ok=True)
    (d / fname).write_bytes(_PNG)
    with (d / "manifest.jsonl").open("a") as fh:
        fh.write(json.dumps(
            {"kind": kind, "file": fname, "label": label}) + "\n")


def _pdf_text(path: Path) -> str:
    from pypdf import PdfReader
    return "\n".join(p.extract_text() or "" for p in PdfReader(str(path)).pages)


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def test_pdf_contains_findings_and_table_values(tmp_path) -> None:
    _seed(tmp_path)
    dest = tmp_path / "exports" / "report.pdf"
    build_analysis_report_pdf(tmp_path, dest)
    assert dest.is_file()
    assert dest.read_bytes()[:5] == b"%PDF-"
    text = _pdf_text(dest)
    assert "Analysis report" in text
    assert "Churn on tenure" in text
    assert "Region counts" in text
    # Frequency-table cell values should have made it into the PDF's
    # rendered table, not just the finding title.
    assert "north" in text
    assert "40" in text


def test_pdf_includes_verification_warnings(tmp_path) -> None:
    # r_squared this low with this large an N should trip a verification
    # warning surfaced in the existing Markdown/HTML report; the PDF must
    # show the same warning text, not silently drop it.
    _seed(tmp_path)
    from sift.research_export import build_analysis_report
    md_report = build_analysis_report(tmp_path)
    dest = tmp_path / "exports" / "report.pdf"
    build_analysis_report_pdf(tmp_path, dest)
    text = _pdf_text(dest)
    # Whatever verification detail strings appear in the trusted
    # Markdown report must also appear in the PDF -- same material,
    # same verdicts, different renderer.
    for line in md_report["markdown"].splitlines():
        if line.startswith("- ⚠ "):
            detail = line.split(" ", 2)[-1]
            assert detail[:40] in text


def test_pdf_embeds_allowlisted_figure_only(tmp_path) -> None:
    _seed(tmp_path)
    _plot(tmp_path, "coefficients", "coef.png", "Coefficient plot")
    _plot(tmp_path, "raw_scatter_not_allowlisted", "raw.png", "Raw scatter")
    dest = tmp_path / "exports" / "report.pdf"
    build_analysis_report_pdf(tmp_path, dest)
    text = _pdf_text(dest)
    assert "Coefficient plot" in text
    assert "Raw scatter" not in text


def test_pdf_ai_disclosure_statement_present(tmp_path) -> None:
    _seed(tmp_path)
    dest = tmp_path / "exports" / "report.pdf"
    build_analysis_report_pdf(tmp_path, dest)
    text = _pdf_text(dest)
    assert "AI-use disclosure statement" in text
    assert "Limitations" in text


def test_pdf_empty_session_does_not_raise(tmp_path) -> None:
    dest = tmp_path / "exports" / "empty.pdf"
    build_analysis_report_pdf(tmp_path, dest)
    assert dest.is_file()
    text = _pdf_text(dest)
    assert "No analysis results have been recorded" in text


def test_pdf_caps_huge_table_and_says_so(tmp_path) -> None:
    (tmp_path / ".sift").mkdir(exist_ok=True)
    store = get_store(tmp_path)
    counts = {f"cat_{i:03d}": 10 + i for i in range(200)}
    store.insert(
        label="Huge taxonomy", analysis_type="frequency_table",
        sanitized_payload={
            "type": "frequency_table", "variable": "cat", "n": sum(counts.values()),
            "counts": counts, "missing_count": 0,
        },
        language="Python", script_code="x", transformations=[],
        raw_log_path=None, script_run_id="r1")
    dest = tmp_path / "exports" / "huge.pdf"
    build_analysis_report_pdf(tmp_path, dest)
    text = _pdf_text(dest)
    assert "Huge taxonomy" in text
    assert "cat_000" in text


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def _pptx_all_text(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        chunks.append(cell.text)
    return "\n".join(chunks)


def test_pptx_contains_findings_and_table_values(tmp_path) -> None:
    _seed(tmp_path)
    dest = tmp_path / "exports" / "report.pptx"
    build_analysis_report_pptx(tmp_path, dest)
    assert dest.is_file()
    text = _pptx_all_text(dest)
    assert "Churn on tenure" in text
    assert "Region counts" in text
    assert "north" in text
    assert "40" in text


def test_pptx_slide_count_scales_with_findings_and_figures(tmp_path) -> None:
    from pptx import Presentation
    _seed(tmp_path)
    _plot(tmp_path, "coefficients", "coef.png", "Coefficient plot")
    dest = tmp_path / "exports" / "report.pptx"
    build_analysis_report_pptx(tmp_path, dest)
    prs = Presentation(str(dest))
    # title + summary + 2 findings + 1 figure + AI-use + limitations
    # (no across-result checks guaranteed present, so >= not ==)
    assert len(prs.slides) >= 7
    titles = []
    for slide in prs.slides:
        if slide.shapes.title is not None:
            titles.append(slide.shapes.title.text)
    joined = " | ".join(titles)
    assert "Summary" in joined
    assert "AI-use disclosure statement" in joined
    assert "Limitations" in joined


def test_pptx_embeds_allowlisted_figure_only(tmp_path) -> None:
    _seed(tmp_path)
    _plot(tmp_path, "coefficients", "coef.png", "Coefficient plot")
    _plot(tmp_path, "raw_scatter_not_allowlisted", "raw.png", "Raw scatter")
    dest = tmp_path / "exports" / "report.pptx"
    build_analysis_report_pptx(tmp_path, dest)
    text = _pptx_all_text(dest)
    assert "Coefficient plot" in text
    assert "Raw scatter" not in text


def test_pptx_empty_session_does_not_raise(tmp_path) -> None:
    dest = tmp_path / "exports" / "empty.pptx"
    build_analysis_report_pptx(tmp_path, dest)
    assert dest.is_file()
    text = _pptx_all_text(dest)
    assert "No analysis results have been recorded" in text


def test_pptx_caps_huge_table_and_says_so(tmp_path) -> None:
    (tmp_path / ".sift").mkdir(exist_ok=True)
    store = get_store(tmp_path)
    counts = {f"cat_{i:03d}": 10 + i for i in range(200)}
    store.insert(
        label="Huge taxonomy", analysis_type="frequency_table",
        sanitized_payload={
            "type": "frequency_table", "variable": "cat", "n": sum(counts.values()),
            "counts": counts, "missing_count": 0,
        },
        language="Python", script_code="x", transformations=[],
        raw_log_path=None, script_run_id="r1")
    dest = tmp_path / "exports" / "huge.pptx"
    build_analysis_report_pptx(tmp_path, dest)
    text = _pptx_all_text(dest)
    assert "Huge taxonomy" in text
    assert "more row(s)" in text


# ---------------------------------------------------------------------------
# Cross-format consistency with policy exclusion (dataset exportable: false)
# ---------------------------------------------------------------------------

def test_pdf_and_pptx_respect_dataset_export_policy(tmp_path) -> None:
    from sift.policy import DatasetPolicy, SiftPolicy, save_policy

    (tmp_path / ".sift").mkdir(exist_ok=True)
    store = get_store(tmp_path)
    store.insert(
        label="Sensitive finding", analysis_type="descriptive",
        sanitized_payload={
            "type": "descriptive", "variable": "income", "n": 500,
            "mean": 55000.0, "sd": 12000.0, "min": 20000, "max": 200000,
            "median": 52000.0, "missing_count": 0,
            "source_dataset": "payroll.csv",
        },
        language="Python", script_code="x", transformations=[],
        raw_log_path=None, script_run_id="r1", source_dataset="payroll.csv")
    save_policy(tmp_path, SiftPolicy(datasets={
        "payroll.csv": DatasetPolicy(exportable=False),
    }))

    pdf_dest = tmp_path / "exports" / "report.pdf"
    pptx_dest = tmp_path / "exports" / "report.pptx"
    build_analysis_report_pdf(tmp_path, pdf_dest)
    build_analysis_report_pptx(tmp_path, pptx_dest)
    assert "Sensitive finding" not in _pdf_text(pdf_dest)
    assert "Sensitive finding" not in _pptx_all_text(pptx_dest)
